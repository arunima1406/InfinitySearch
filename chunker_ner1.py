import pdfplumber
import fitz
import re
import json
import glob
import os
import time
import uuid
import google.generativeai as genai
from PIL import Image
from dotenv import load_dotenv
import docx
import csv
import openpyxl
from pptx import Presentation
import boto3

# --- Configuration ---
IMAGE_OUTPUT_DIR = "/tmp/extracted_images"  # Cloud-friendly temp folder
MAX_TEXT_LENGTH_FOR_SUMMARY = 200000
QUICK_TASK_MODEL = 'gemini-2.5-flash-lite'
EXTRACTION_MODEL = 'gemini-2.5-flash-lite'
API_DELAY_SECONDS = 3
BUCKET_NAME = "infinity-search-v1"
S3_FOLDER = "kg-json"
# Directory to save intermediate chunk files for debugging
INTERMEDIATE_CHUNKS_DIR = "intermediate_chunks"
# Set to True to save the chunks for each file as a separate JSON
SAVE_INTERMEDIATE_CHUNKS = True
# Parameters for text chunking
CHUNK_SIZE = 250  # Increased for better context
OVERLAP_SIZE = 50

# --- Setup ---
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise ValueError("API key not found. Please set GOOGLE_API_KEY in your .env file.")
genai.configure(api_key=api_key)

s3 = boto3.client(
    "s3",
    aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
    aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY"),
    region_name="ap-south-1"
)

# --- Gemini Models and Prompts ---
quick_model = genai.GenerativeModel(QUICK_TASK_MODEL)
extraction_model = genai.GenerativeModel(
    EXTRACTION_MODEL,
    generation_config={"response_mime_type": "application/json"}
)
TRIPLE_EXTRACTION_PROMPT = """
You are an expert information extraction agent. Your task is to analyze the user's text and extract knowledge triples.

Output ONLY a valid JSON array of objects. Each object represents a single relationship and must have three keys: "start", "relation", and "end".

- The "relation" value must be descriptive, uppercase, and snake_case (e.g., PROPOSED, IS_A).
- Do not add any explanations or introductory text outside of the main JSON array.

Example for "Alan Turing proposed the universal machine in 1936.":
[
  {
    "start": "Alan Turing",
    "relation": "PROPOSED",
    "end": "universal machine"
  },
  {
    "start": "universal machine",
    "relation": "PROPOSED_IN",
    "end": "1936"
  }
]
"""

# --- Helper Functions ---
def describe_image_with_gemini(image_path):
    print(f"      > Describing image: {os.path.basename(image_path)}...")
    try:
        img = Image.open(image_path)
        prompt = "Analyze this image in detail. Provide a comprehensive, descriptive paragraph explaining its content. Describe the main subject, the setting, any actions taking place, colors, textures, and the overall mood or context. Be as thorough as possible, as if describing it to someone who cannot see it. Discard the use of wasteful words, and keep it precise but descriptive. always state the subject, dont use pronouns, and keep in mind this for triples extraction."
        response = quick_model.generate_content([prompt, img])
        time.sleep(API_DELAY_SECONDS)
        return response.text.strip()
    except Exception as e:
        print(f"        ! Error describing image: {e}")
        return "Error during image description."

def summarize_text_with_gemini(text):
    print("    > Summarizing document text...")
    if len(text) > MAX_TEXT_LENGTH_FOR_SUMMARY:
        text = text[:MAX_TEXT_LENGTH_FOR_SUMMARY]
        print("      ! Text was truncated for summarization due to extreme length.")
    try:
        prompt = "Analyze the following. Synthesize its core subject and key points into a single, concise, and factual paragraph. Describe the content directly as if explaining it, starting with the main subject. Do not use third-person language like 'This text is about' or 'The author discusses'."
        response = quick_model.generate_content([prompt, text])
        time.sleep(API_DELAY_SECONDS)
        return response.text.strip()
    except Exception as e:
        print(f"      ! Error during text summarization: {e}")
        return "Error during summarization."

def create_overlapping_chunks(text, chunk_size, overlap_size):
    words = text.split()
    if len(words) <= chunk_size:
        return [" ".join(words)]
    chunks = []
    step = chunk_size - overlap_size
    for i in range(0, len(words), step):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks

# --- File Processors ---
def process_pdf(file_path):
    print("    > Extracting text and images from PDF...")
    full_text, images_data = "", []
    with pdfplumber.open(file_path) as pdf:
        full_text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    try:
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            for j, img in enumerate(page.get_images(full=True), start=1):
                xref, base_image = img[0], doc.extract_image(img[0])
                image_bytes, image_ext = base_image["image"], base_image["ext"]
                img_filename = f"{os.path.basename(file_path)}_p{i+1}_i{j}.{image_ext}"
                img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
                with open(img_path, "wb") as f: f.write(image_bytes)
                description = describe_image_with_gemini(img_path)
                images_data.append({"image_path": img_path, "description": description})
        doc.close()
    except Exception as e:
        print(f"      ! Could not process images for PDF: {e}")
    return full_text, images_data

def process_docx(file_path):
    print("    > Extracting text and images from DOCX...")
    doc = docx.Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs])
    images_data = []
    for i, rel in enumerate(doc.part.rels):
        if "image" in doc.part.rels[rel].target_ref:
            img_data, img_ext = doc.part.rels[rel].target_part.blob, doc.part.rels[rel].target_ref.split('.')[-1]
            img_filename = f"{os.path.basename(file_path)}_i{i+1}.{img_ext}"
            img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
            with open(img_path, "wb") as f: f.write(img_data)
            description = describe_image_with_gemini(img_path)
            images_data.append({"image_path": img_path, "description": description})
    return full_text, images_data

def process_txt(file_path):
    print("    > Extracting text from TXT...")
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read(), []

def process_pptx(file_path):
    print("    > Extracting text and images from PPTX...")
    prs = Presentation(file_path)
    full_text = ""
    images_data = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
            if shape.shape_type == 13:  # Picture
                img = shape.image
                img_bytes = img.blob
                img_ext = img.ext
                img_filename = f"{os.path.basename(file_path)}_s{slide_number}_i{len(images_data)+1}.{img_ext}"
                img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                description = describe_image_with_gemini(img_path)
                images_data.append({"image_path": img_path, "description": description})
    return full_text, images_data

def process_csv(file_path):
    print("    > Extracting text from CSV...")
    full_text = ""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            full_text += " ".join(row) + "\n"
    return full_text, []

def process_image(file_path):
    print("    > Processing image file...")
    description = describe_image_with_gemini(file_path)
    return "", [{"image_path": file_path, "description": description}]

# --- Knowledge Extraction ---
def extract_knowledge_triples(texts_to_process):
    print("    > Extracting knowledge triples...")
    all_triples = []
    for i, text_chunk in enumerate(texts_to_process, start=1):
        print(f"      > Processing text chunk {i}/{len(texts_to_process)}...")
        prompt_parts = [TRIPLE_EXTRACTION_PROMPT, "Text to analyze:", text_chunk]
        try:
            response = extraction_model.generate_content(prompt_parts)
            # Clean the response
            clean_response = response.text.strip()
            if clean_response.startswith('```json'):
                clean_response = clean_response[7:]
            if clean_response.endswith('```'):
                clean_response = clean_response[:-3]
            clean_response = clean_response.strip()
            list_of_dicts = json.loads(clean_response)
            print(f"        + Extracted {len(list_of_dicts)} triples.")
            all_triples.extend(list_of_dicts)
            time.sleep(API_DELAY_SECONDS)
        except Exception as e:
            print(f"        ! Error processing a text part: {e}")
    return all_triples

# --- Single File Processing ---
def process_single_file(file_path, user_id):
    print(f"--- 🎬 Processing File: {file_path} for User ID: {user_id} ---")
    full_text, images_data = "", []
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf": full_text, images_data = process_pdf(file_path)
        elif ext == ".docx": full_text, images_data = process_docx(file_path)
        elif ext == ".txt": full_text, images_data = process_txt(file_path)
        elif ext == ".pptx": full_text, images_data = process_pptx(file_path)
        elif ext == ".csv": full_text, images_data = process_csv(file_path)
        elif ext in [".png", ".jpg", ".jpeg"]: full_text, images_data = process_image(file_path)

        cleaned_text = re.sub(r'\s+', ' ', full_text).strip()
        summary = summarize_text_with_gemini(cleaned_text) if cleaned_text else "No text to summarize."

        all_content_text = cleaned_text + " " + " ".join([img["description"] for img in images_data])
        text_chunks = create_overlapping_chunks(all_content_text, CHUNK_SIZE, OVERLAP_SIZE)

        if SAVE_INTERMEDIATE_CHUNKS:
            os.makedirs(INTERMEDIATE_CHUNKS_DIR, exist_ok=True)
            chunk_filename = f"chunks_{os.path.basename(file_path)}.json"
            chunk_output_path = os.path.join(INTERMEDIATE_CHUNKS_DIR, chunk_filename)
            with open(chunk_output_path, "w", encoding="utf-8") as f:
                json.dump(text_chunks, f, indent=4)
            print(f"    > Intermediate chunks saved to: {chunk_output_path}")

        extracted_triples = extract_knowledge_triples(text_chunks)

        unique_triples_list, seen_triples = [], set()
        for triple in extracted_triples:
            canonical_triple = (
                triple.get('start', '').strip().lower(),
                triple.get('relation', '').strip().lower(),
                triple.get('end', '').strip().lower()
            )
            if all(canonical_triple) and canonical_triple not in seen_triples:
                seen_triples.add(canonical_triple)
                unique_triples_list.append(triple)

        triples_dict = {str(i): triple for i, triple in enumerate(unique_triples_list, 1)}

        file_knowledge = {
            "user_id": user_id,
            "source_file": os.path.basename(file_path),
            "episode_id": f"ep_{uuid.uuid4()}",
            "summary": summary,
            "triples": triples_dict
        }

        print(f"--- ✅ Finished Processing: {file_path} ---\n")
        return file_knowledge
    except Exception as e:
        print(f"--- ❌ FAILED to process {file_path}. Error: {e} ---\n")
        return None

# --- Main Execution Pipeline ---
def main():
    test_user_id = "123456789-local-test-user"
    os.makedirs(IMAGE_OUTPUT_DIR, exist_ok=True)
    if SAVE_INTERMEDIATE_CHUNKS:
        os.makedirs(INTERMEDIATE_CHUNKS_DIR, exist_ok=True)
    file_types = ["*.pdf", "*.docx", "*.txt", "*.pptx", "*.csv", "*.png", "*.jpg", "*.jpeg"]
    all_files = [file for ftype in file_types for file in glob.glob(ftype)]

    if not all_files:
        print("No files found to process.")
        return

    final_knowledge_graph_data = []
    for file_path in all_files:
        processed_data = process_single_file(file_path, test_user_id)
        if processed_data:
            final_knowledge_graph_data.append(processed_data)

    output_filename = "knowledge_graph_with_userid.json"
    with open(output_filename, "w", encoding="utf-8") as f:
        json.dump({"knowledge_graph": final_knowledge_graph_data}, f, indent=4, ensure_ascii=False)

    print(f"Knowledge graph saved to '{output_filename}'")

    # --- Upload to S3 ---
    try:
        s3_object_key = f"{S3_FOLDER}/{output_filename}"
        s3.upload_file(output_filename, BUCKET_NAME, s3_object_key)
        public_url = f"https://{BUCKET_NAME}.s3.amazonaws.com/{s3_object_key}"
        print(f"✅ JSON uploaded successfully to S3!")
        print(f"🌐 File URL: {public_url}")
    except Exception as e:
        print(f"❌ Failed to upload JSON to S3: {e}")
if __name__ == "__main__":
    main()
