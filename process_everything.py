import pdfplumber
import fitz
import re
import json
import glob
import os
import google.generativeai as genai
from PIL import Image
from dotenv import load_dotenv
import docx
import csv
import openpyxl
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor

# --- Configuration ---
IMAGE_OUTPUT_DIR = "extracted_images"
CHUNK_SIZE = 200
OVERLAP_SIZE = 40
MAX_WORKERS = 5 # Number of concurrent API calls. Adjust carefully!
# ---------------------

# Load API Key and configure Gemini
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# --- Helper Functions (describe_image and create_chunks are the same) ---

def describe_image_with_gemini(image_path):
    # This function is the "task" our assistants will do.
    # It remains the same, but the print statement is removed to avoid jumbled output.
    try:
        img = Image.open(image_path)
        model = genai.GenerativeModel('gemini-1.5-flash-latest')
        response = model.generate_content(["Describe this image in a concise sentence.", img])
        return response.text
    except Exception as e:
        # It's useful to know which file failed in a concurrent run
        print(f"      ! Error describing image '{os.path.basename(image_path)}': {e}")
        return "Error during Gemini description."

def create_overlapping_chunks(text, chunk_size, overlap_size):
    words = text.split()
    if len(words) <= chunk_size: return [" ".join(words)]
    chunks, step = [], chunk_size - overlap_size
    for i in range(0, len(words), step):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks

# --- NEW Concurrency Helper ---
def describe_images_concurrently(image_paths):
    descriptions = []
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        # The map function sends each image_path to the description function
        # and collects the results in the same order.
        results = executor.map(describe_image_with_gemini, image_paths)
        descriptions = list(results)
    return descriptions

# --- Updated File Processing Functions ---

def process_pdf(file_path):
    full_text, images_data = "", []
    with pdfplumber.open(file_path) as pdf:
        full_text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    
    image_paths_to_describe = []
    try:
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            for j, img in enumerate(page.get_images(full=True), start=1):
                xref, base_image = img[0], doc.extract_image(img[0])
                image_bytes, image_ext = base_image["image"], base_image["ext"]
                img_filename = f"{os.path.basename(file_path)}_p{i+1}_i{j}.{image_ext}"
                img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
                with open(img_path, "wb") as f: f.write(image_bytes)
                image_paths_to_describe.append(img_path)
        doc.close()
        
        print(f"    > Found {len(image_paths_to_describe)} images. Describing them concurrently...")
        descriptions = describe_images_concurrently(image_paths_to_describe)
        images_data = [{"image_path": path, "description": desc} for path, desc in zip(image_paths_to_describe, descriptions)]

    except Exception as e: print(f"    ! Could not process images for PDF: {e}")
    return full_text, images_data

def process_docx(file_path):
    doc = docx.Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs])
    image_paths_to_describe = []
    for i, rel in enumerate(doc.part.rels):
        if "image" in doc.part.rels[rel].target_ref:
            img_data, img_ext = doc.part.rels[rel].target_part.blob, doc.part.rels[rel].target_ref.split('.')[-1]
            img_filename = f"{os.path.basename(file_path)}_i{i+1}.{img_ext}"
            img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
            with open(img_path, "wb") as f: f.write(img_data)
            image_paths_to_describe.append(img_path)
    
    print(f"    > Found {len(image_paths_to_describe)} images. Describing them concurrently...")
    descriptions = describe_images_concurrently(image_paths_to_describe)
    images_data = [{"image_path": path, "description": desc} for path, desc in zip(image_paths_to_describe, descriptions)]
    return full_text, images_data

def process_pptx(file_path):
    prs, full_text, image_paths_to_describe = Presentation(file_path), "", []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text'): full_text += shape.text + "\n"
            if hasattr(shape, 'image'):
                img, img_ext = shape.image, shape.image.ext
                img_filename = f"{os.path.basename(file_path)}_s{i+1}.{img_ext}"
                img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
                with open(img_path, "wb") as f: f.write(img.blob)
                image_paths_to_describe.append(img_path)

    print(f"    > Found {len(image_paths_to_describe)} images. Describing them concurrently...")
    descriptions = describe_images_concurrently(image_paths_to_describe)
    images_data = [{"image_path": path, "description": desc} for path, desc in zip(image_paths_to_describe, descriptions)]
    return full_text, images_data

# The text-only processors remain the same
def process_xlsx(file_path):
    workbook = openpyxl.load_workbook(file_path)
    full_text = ""
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            full_text += " ".join([str(cell.value) for cell in row if cell.value is not None]) + "\n"
    return full_text, []
def process_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f: return f.read(), []
def process_csv(file_path):
    text = ""
    with open(file_path, 'r', encoding='utf-8') as f:
        for row in csv.reader(f): text += " ".join(row) + "\n"
    return text, []

# --- Main Logic (Mostly the same) ---
# ... (Code from previous step to find all files and loop through them) ...
os.makedirs(IMAGE_OUTPUT_DIR, exist_ok=True)
file_types = ["*.pdf", "*.docx", "*.txt", "*.csv", "*.xlsx", "*.pptx"]
all_files = [file for ftype in file_types for file in glob.glob(ftype)]
consolidated_results = []
print(f"Found {len(all_files)} file(s) to process.")

for file_path in all_files:
    print(f"--- Processing: {file_path} ---")
    full_text, images_data = "", []
    try:
        ext = file_path.lower().split('.')[-1]
        if ext == "pdf": full_text, images_data = process_pdf(file_path)
        elif ext == "docx": full_text, images_data = process_docx(file_path)
        elif ext == "pptx": full_text, images_data = process_pptx(file_path)
        elif ext == "xlsx": full_text, images_data = process_xlsx(file_path)
        elif ext == "txt": full_text, images_data = process_txt(file_path)
        elif ext == "csv": full_text, images_data = process_csv(file_path)

        cleaned_text = re.sub(r'[^a-zA-Z0-9\s]', '', full_text.lower())
        text_chunks = create_overlapping_chunks(cleaned_text, CHUNK_SIZE, OVERLAP_SIZE)
        
        file_data = {"source_file": file_path, "text_chunks": text_chunks, "extracted_images": images_data}
        consolidated_results.append(file_data)
        
    except Exception as e: print(f"    ! FAILED to process {file_path}. Error: {e}")

output_filename = "final_document_output.json"
with open(output_filename, "w") as f: json.dump(consolidated_results, f, indent=4)

print(f"\n--- ALL DONE! --- \nComprehensive data saved to '{output_filename}'")