import pdfplumber
import fitz
import re
import json
import glob
import os
import time
import uuid
import google.generativeai as genai
from PIL import Image
from dotenv import load_dotenv
import docx
import csv  # NEW: Import for CSV handling
from pptx import Presentation  # NEW: Import for PowerPoint handling

# --- Configuration ---
IMAGE_OUTPUT_DIR = "/tmp/extracted_images"  # Cloud-friendly temp folder
MAX_TEXT_LENGTH_FOR_SUMMARY = 200000
QUICK_TASK_MODEL = 'gemini-2.5-flash-lite'
EXTRACTION_MODEL = 'gemini-2.5-flash-lite'
API_DELAY_SECONDS = 3

# --- Setup, Models, and Prompts (Unchanged from previous version) ---
# (The setup code for API keys, models, and prompts remains the same)
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise ValueError("API key not found. Please set GOOGLE_API_KEY in your .env file.")
genai.configure(api_key=api_key)

quick_model = genai.GenerativeModel(QUICK_TASK_MODEL)
extraction_model = genai.GenerativeModel(
    EXTRACTION_MODEL,
    generation_config={"response_mime_type": "application/json"}
)
TRIPLE_EXTRACTION_PROMPT = """
You are an expert information extraction agent. Your task is to analyze the user's text and extract knowledge triples.
Output ONLY a valid JSON array of objects. Each object represents a single relationship and must have three keys: "start", "relation", and "end".
- The "relation" value must be descriptive, uppercase, and snake_case (e.g., DEPICTS, IS_A, HAS_PROPERTY).
- Do not add any explanations or introductory text outside of the main JSON array.
"""
# --- Helper Functions (Unchanged from previous version) ---
def describe_image_with_gemini(image_path):
    print(f"        > Describing image: {os.path.basename(image_path)}...")
    try:
        img = Image.open(image_path)
        prompt = "Analyze this image in detail. Provide a comprehensive, descriptive paragraph explaining its content. Describe the main subject, the setting, any actions taking place, colors, textures, and the overall mood or context. Be as thorough as possible, as if describing it to someone who cannot see it. Discard the use of wasteful words, and keep it precise but descriptive. always state the subject, dont use pronouns, and keep in mind this for triples extraction."
        response = quick_model.generate_content([prompt, img])
        time.sleep(API_DELAY_SECONDS)
        return response.text.strip()
    except Exception as e:
        print(f"        ! Error describing image: {e}")
        return "Error during image description."

def summarize_text_with_gemini(text):
    print("    > Summarizing document text...")
    if len(text) > MAX_TEXT_LENGTH_FOR_SUMMARY:
        text = text[:MAX_TEXT_LENGTH_FOR_SUMMARY]
    try:
        prompt = "Analyze the following. Synthesize its core subject and key points into a single, concise, and factual paragraph. Describe the content directly as if explaining it, starting with the main subject. Do not use third-person language like 'This text is about' or 'The author discusses'."
        response = quick_model.generate_content([prompt, text])
        time.sleep(API_DELAY_SECONDS)
        return response.text.strip()
    except Exception as e:
        print(f"      ! Error during text summarization: {e}")
        return "Error during summarization."

# --- File Processors ---
def process_pdf(file_path):
    # (This function is unchanged)
    print("    > Extracting text and images from PDF...")
    full_text, images_data = "", []
    with pdfplumber.open(file_path) as pdf:
        full_text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    try:
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            for j, img in enumerate(page.get_images(full=True), start=1):
                xref, base_image = img[0], doc.extract_image(img[0])
                image_bytes, image_ext = base_image["image"], base_image["ext"]
                img_filename = f"{os.path.splitext(os.path.basename(file_path))[0]}_p{i+1}_i{j}.{image_ext}"
                img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
                with open(img_path, "wb") as f: f.write(image_bytes)
                description = describe_image_with_gemini(img_path)
                images_data.append({"image_path": img_path, "description": description})
        doc.close()
    except Exception as e:
        print(f"      ! Could not process images for PDF: {e}")
    return full_text, images_data

def process_docx(file_path):
    # (This function is unchanged)
    print("    > Extracting text and images from DOCX...")
    doc = docx.Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs])
    images_data = []
    for i, rel in enumerate(doc.part.rels):
        if "image" in doc.part.rels[rel].target_ref:
            img_data = doc.part.rels[rel].target_part.blob
            img_ext = doc.part.rels[rel].target_ref.split('.')[-1]
            img_filename = f"{os.path.splitext(os.path.basename(file_path))[0]}_i{i+1}.{img_ext}"
            img_path = os.path.join(IMAGE_OUTPUT_DIR, img_filename)
            with open(img_path, "wb") as f: f.write(img_data)
            description = describe_image_with_gemini(img_path)
            images_data.append({"image_path": img_path, "description": description})
    return full_text, images_data

def process_txt(file_path):
    # (This function is unchanged)
    print("    > Extracting text from TXT...")
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read(), []

def process_image(file_path):
    # (This function is unchanged)
    print("    > Processing standalone image file...")
    description = describe_image_with_gemini(file_path)
    return description, []

# NEW: Processor for PPTX files
def process_pptx(file_path):
    """
    Extracts all text from a PowerPoint .pptx file.
    Note: This version does not extract images from slides.
    """
    print("    > Extracting text from PPTX...")
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs), []
    except Exception as e:
        print(f"      ! Error processing PPTX file: {e}")
        return "", []

# NEW: Processor for CSV files
def process_csv(file_path):
    """
    Reads a CSV file and converts its content into a single string.
    """
    print("    > Extracting text from CSV...")
    try:
        full_text_parts = []
        with open(file_path, mode='r', encoding='utf-8', errors='ignore') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                # Convert each row into a comma-separated string
                full_text_parts.append(", ".join(row))
        # Join all rows with a newline to preserve structure
        return "\n".join(full_text_parts), []
    except Exception as e:
        print(f"      ! Error processing CSV file: {e}")
        return "", []

# --- Knowledge Extraction (Unchanged) ---
def extract_knowledge_triples(full_text):
    print("    > Extracting knowledge triples from content...")
    if not full_text.strip():
        return []
    prompt_parts = [TRIPLE_EXTRACTION_PROMPT, "Text to analyze:", full_text]
    try:
        response = extraction_model.generate_content(prompt_parts)
        cleaned_response = response.text.strip().replace("```json", "").replace("```", "")
        list_of_dicts = json.loads(cleaned_response)
        print(f"      + Extracted {len(list_of_dicts)} triples.")
        time.sleep(API_DELAY_SECONDS)
        return list_of_dicts
    except Exception as e:
        print(f"      ! Error processing content for triples: {e}")
        return []

# --- Single File Processing Pipeline (MODIFIED) ---
def process_single_file(file_path, user_id):
    print(f"--- 🎬 Processing File: {os.path.basename(file_path)} for User ID: {user_id} ---")
    full_text, images_data = "", []
    image_extensions = ['.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif']
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf": full_text, images_data = process_pdf(file_path)
        elif ext == ".docx": full_text, images_data = process_docx(file_path)
        elif ext == ".txt": full_text, images_data = process_txt(file_path)
        elif ext in image_extensions: full_text, images_data = process_image(file_path)
        elif ext == ".pptx": full_text, images_data = process_pptx(file_path) # MODIFIED
        elif ext == ".csv": full_text, images_data = process_csv(file_path) # MODIFIED
        else:
             print(f"   ! Skipping unsupported file type: {ext}")
             return None

        cleaned_text = re.sub(r'\s+', ' ', full_text).strip()
        summary = cleaned_text if ext in image_extensions else summarize_text_with_gemini(cleaned_text)
        if not cleaned_text:
            summary = "No text to summarize."

        all_content_text = cleaned_text + " " + " ".join([img["description"] for img in images_data])
        extracted_triples = extract_knowledge_triples(all_content_text)

        unique_triples_list, seen_triples = [], set()
        for triple in extracted_triples:
            canonical_triple = (
                str(triple.get('start', '')).strip().lower(),
                str(triple.get('relation', '')).strip().lower(),
                str(triple.get('end', '')).strip().lower()
            )
            if all(canonical_triple) and canonical_triple not in seen_triples:
                seen_triples.add(canonical_triple)
                unique_triples_list.append(triple)

        triples_dict = {str(i): triple for i, triple in enumerate(unique_triples_list, 1)}

        file_knowledge = {
            "user_id": user_id,
            "source_file": os.path.basename(file_path),
            "episode_id": f"ep_{uuid.uuid4()}",
            "summary": summary,
            "triples": triples_dict
        }

        print(f"--- ✅ Finished Processing: {os.path.basename(file_path)} ---\n")
        return file_knowledge
    except Exception as e:
        print(f"--- ❌ FAILED to process {file_path}. Error: {e} ---\n")
        return None

# --- Main Execution (MODIFIED) ---
def main():
    test_user_id = "123456789-local-test-user"
    os.makedirs(IMAGE_OUTPUT_DIR, exist_ok=True)
    
    # MODIFIED: Added PPTX and CSV to the list of file types
    file_types = ["*.pdf", "*.docx", "*.txt", "*.png", "*.jpg", "*.jpeg", "*.webp", "*.pptx", "*.csv"]
    all_files = [file for ftype in file_types for file in glob.glob(ftype)]

    if not all_files:
        print("No supported files found to process.")
        return

    print(f"Found {len(all_files)} files to process.\n")
    final_knowledge_graph_data = []
    for file_path in all_files:
        processed_data = process_single_file(file_path, test_user_id)
        if processed_data:
            final_knowledge_graph_data.append(processed_data)

    if not final_knowledge_graph_data:
        print("No data was successfully processed. Exiting.")
        return

    output_filename = "knowledge_graph_output.json"
    with open(output_filename, "w", encoding="utf-8") as f:
        json.dump({"knowledge_graph": final_knowledge_graph_data}, f, indent=4, ensure_ascii=False)

    print(f"🎉 Knowledge graph processing complete. Output saved to '{output_filename}'")
    # S3 upload logic would go here if enabled...

if __name__ == "__main__":
    main()